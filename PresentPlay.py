from pptx import Presentation
from pptx.util import Inches, Pt
import numpy as np
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT as PP_ALIGN
import math
import sklearn
from sklearn import cluster
from pptx.dml.color import RGBColor, _NoneColor
from pptx.enum.dml import MSO_COLOR_TYPE
import Playgrounds as pg
import matplotlib.pyplot as plt
import matplotlib.pylab as pylab
from operator import itemgetter
import random as rand
from numpy import array
from scipy.cluster.vq import vq, kmeans, whiten


K = 5
UPPERBOUND = 999999999999
LOOK_DISTANCE = 4000000
EMPTY = np.array([1.0, 1.0, 1.0, 1.0, 1.0, 1.0, 1.0, 1.0, 1.0, 1.0, 1.0, 1.0, 1.0, 1.0, 1.0, 1.0, 1.0, 1.0])
EMPTY2 = np.array([0.0, 0.0, 0.0, 0.0, 0.0, 0.0, 0.0, 0.0, 0.0, 0.0, 0.0, 0.0, 0.0, 0.0, 0.0, 0.0, 0.0, 0.0])
HEAD = np.array([0, 0, 0, 1, 0.0000001, 10000000, 5000000, 4500000, 0.000002, 3, 1, 0, -2, 1, -5, 0, 0, 0])
SUBHEAD = np.array([0, 0, 0, 1, 0.00000007, 10000000, 4000000, 3500000, 0.000001, 1, 2, 1, -3, 3, 6, 0, 0, 0])
SUBSECTION = np.array([0, 0, 0, 3, 0.00000001, 13000000, 0.00000001, 3500000, 0.000001, 0, 3, 1, -5, 0, -5, 0, 0, 0])
BODY = np.array([0, 0, 0, 0.01, 0.00000004, 0.000002, 0.00000001, 0.00000001, 1, -3, -2, -1, 4, 0, -50, 0, 0, 0])
CAPTION = np.array([0, 0, 0, 0.007, 0.00000002, 7000000, 0.00000001, 0.00000001, 300000, -2, -1, 3, 2, 0, -5, 0, 0, 0])
WEIGHT_MATRIX = np.array([HEAD, SUBHEAD, SUBSECTION, BODY, CAPTION])
CHANGE_MATRIX = np.array([EMPTY, EMPTY, EMPTY, EMPTY, EMPTY])
STEP = 0.01

prs = Presentation()
title_slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
extraneous = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(1), Inches(1))
extra = extraneous.text_frame
ex = extra.add_paragraph()

title.text = "Hello, World!"
subtitle.text = "this was a triumph!"
ex.text = "It's, like, autotune scat."
ex.font.size = Pt(12)

prs.save('test.pptx')

prs2 = Presentation('thissatest copy.pptx')
#prs1 = Presentation('sectest.pptx')
prs1 = Presentation('leepostertest copy.pptx')
first = prs1.slides[0]
second = prs2.slides[0]
title = first.shapes[0]
#title.text = "Hey, this about to get raunchy"


# Let's assemble a very rough max finding function
titl_tx = ex


def last_lett(word):
    i = -1
    while i > -len(word):
        if not word[i].isspace():
            return word[i]
        i -= 1


def sigmoid(x):
    return 1/(1 + np.exp(-x))


def dsigmoid(x):
    return np.exp(-x)/((1 + np.exp(-x)) ** 2)


def vector_fill_clean(first_slide, run_external):
    """Fills out the np array of all text runs in the document
    Attributes of each vector are in the order:
    1)Run ID 2)Shape ID 3)"""
    """Keeps tabs on the id number of each run and retains them in an array"""
    id_par = 0
    id_run = 0
    vector_arr = []
    del run_external[:]
    for shape in first_slide.shapes:
        if shape.has_text_frame:
            id_shp = first_slide.shapes.index(shape)
            for par in shape.text_frame.paragraphs:
                print(par.text)
                print("This is paragraph: " + str(id_par))
                print("*********")
                for run in par.runs:
                    correct_id = -1
                    if run.font.color.type is not None:
                        if run.font.color.type is MSO_COLOR_TYPE.RGB:
                            if run.font.color.rgb == RGBColor.from_string("FA6300"):
                                correct_id = 0
                            elif run.font.color.rgb == RGBColor.from_string("1398FF"):
                                correct_id = 1
                            elif run.font.color.rgb == RGBColor.from_string("13BA33"):
                                correct_id = 2
                            elif run.font.color.rgb == RGBColor.from_string("000000"):
                                correct_id = 3
                            elif run.font.color.rgb == RGBColor.from_string("EB1500"):
                                correct_id = 4
                    print("The real answer is " + str(correct_id))
                    f_black = 0
                    f_size = 0
                    if (run.font.name is None) or (len(run.font.name) < 5):
                        f_black = 0
                    else:
                        f_black = int(run.font.name[-5:] in ("Black", "Heavy"))
                    if run.font.size is None:
                        f_size = 0
                    else:
                        f_size = float(run.font.size.pt)
                    f_bold = int(run.font.bold == True)
                    f_italic = int(run.font.italic == True)
                    width = shape.width // 1000000
                    height = shape.height // 1000000
                    x = (shape.left // 1000000) + (width / 2)
                    y = (shape.top // 1000000) + (height / 2)
                    l_space = par.line_spacing if par.line_spacing is not None else 0
                    level = par.level
                    length = len(run.text)
                    punc = (int(last_lett(run.text) in (".", ",", "!", "?"))) if len(run.text) > 0 else 0
                    sect = int(run.text in ('Abstract', 'ABSTRACT', 'Introduction', 'INTRODUCTION', 'Method', 'METHOD',
                                            'Results', 'RESULTS', "Acknowledgments", "ACKNOWLEDGEMENTS", "Discussion",
                                            "DISCUSSION", "References", "REFERENCES", "Conclusions", "CONCLUSIONS"))
                    align = int(par.alignment == PP_ALIGN.CENTER)
                    run_attr = [id_run, id_shp, correct_id, length, width, height, x, y, f_size, f_black, f_bold, f_italic, punc,
                                align, sect, l_space, level, 0]
                    vector_arr.append(run_attr)
                    run_external.append(run)
                    id_run += 1
                id_par += 1
    nump_arr = np.array(vector_arr)
    print(nump_arr[40:])
    return nump_arr


def vector_fill_shape(first_slide, numpy_text):
    vector_arr = []
    for shape in first_slide.shapes:
        id = first_slide.shapes.index(shape)
        x = shape.left
        y = shape.top
        #width = shape.width
        #height = shape.height
        tru_x = (shape.width // 2) + x
        tru_y = (shape.height // 2) + y
        text = shape.has_text_frame
        image = int(shape.name[:7] == "Picture")
        chart = shape.has_chart
        table = shape.has_table
        head = subhead = subsec = body = subtitl = 0
        if text == 1:
            for run in numpy_text:
                #Checks that the id of the numpy element is the same as the shape id
                if run[1] == id:
                    head = int(run[-1] == 0)
                    subhead = int(run[-1] == 1)
                    subsec = int(run[-1] == 2)
                    body = int(run[-1] == 3)
                    subtitl = int(run[-1] == 4)

        shape_attr = [id, x, y, tru_x, tru_y, text, image, chart, table, head, subhead, subsec, body, subtitl, 0]
        vector_arr.append(shape_attr)
    return np.array(vector_arr)


def shape_vector_distill(shape_np):
    distilled = []
    for vect in shape_np:
        distilled.append([vect[3], vect[4]])
    return np.array(distilled)


def euclid_distance(x, xi):
    return np.sqrt(np.sum((x - xi)**2))


def neighbourhood_points(X, x_centroid, distance = 5):
    eligible_X = []
    for x in X:
        distance_between = euclid_distance(x, x_centroid)
        # print('Evaluating: [%s vs %s] yield dist=%.2f' % (x, x_centroid, distance_between))
        if distance_between <= distance:
            eligible_X.append(x)
    return eligible_X


def gaussian_kernel(distance, bandwidth):
    val = (1/(bandwidth*math.sqrt(2*math.pi))) * np.exp(-0.5*((distance / bandwidth))**2)
    return val


def meanshift(original_X):
    original_X = shape_vector_distill(original_X)
    kernel_bandwidth = cluster.estimate_bandwidth(original_X)
    X = np.copy(original_X)

    past_X = []
    n_iterations = 45
    print("is this working?")
    for it in range(n_iterations):
        for i, x in enumerate(X):
            # Step 1. For each datapoint x ∈ X, find the neighbouring points N(x) of x.
            print("yep yep")
            neighbours = neighbourhood_points(X, x, LOOK_DISTANCE)

            # Step 2. For each datapoint x ∈ X, calculate the mean shift m(x).
            numerator = 0
            denominator = 0
            for neighbour in neighbours:
                distance = euclid_distance(neighbour, x)
                weight = gaussian_kernel(LOOK_DISTANCE, kernel_bandwidth)
                numerator += (weight * neighbour)
                denominator += weight

            new_x = numerator / denominator

            ### Step 3. For each datapoint x ∈ X, update x ← m(x).
            X[i] = new_x
            print(new_x)

        past_X.append(np.copy(X))
    return X


def arrange_meanshift(meanshifted, shape_numpy):
    seen = []

    i = 0
    while i < len(meanshifted):
        vect = meanshifted[i].tolist()
        if vect in seen:
            shape_numpy[i][-1] = seen.index(vect)
        else:
            seen.append(meanshifted[i].tolist())
            shape_numpy[i][-1] = seen.index(vect)
        i += 1


def distr_assigning(nparray):
    head_size = 88
    total_runs = 0
    correct_runs = 0
    change_matrix = CHANGE_MATRIX
    for vect in nparray:
        # print(vect)
        # print(WEIGHT_MATRIX[0])
        # w_head = (0, sigmoid(np.dot(WEIGHT_MATRIX[0], vect)))
        # w_subhead = (1, sigmoid(np.dot(WEIGHT_MATRIX[1], vect)))
        # w_subsec = (2, sigmoid(np.dot(WEIGHT_MATRIX[2], vect)))
        # w_body = (3, sigmoid(np.dot(WEIGHT_MATRIX[3], vect)))
        # w_subtitle = (4, sigmoid(np.dot(WEIGHT_MATRIX[4], vect)))
        # weights = [w_head, w_subhead, w_subsec, w_body, w_subtitle]

        w_head = (0, (np.dot(WEIGHT_MATRIX[0], vect)))
        w_subhead = (1, (np.dot(WEIGHT_MATRIX[1], vect)))
        w_subsec = (2, (np.dot(WEIGHT_MATRIX[2], vect)))
        w_body = (3, (np.dot(WEIGHT_MATRIX[3], vect)))
        w_subtitle = (4, (np.dot(WEIGHT_MATRIX[4], vect)))
        weights = [w_head, w_subhead, w_subsec, w_body, w_subtitle]

        weights.sort(key=lambda x: x[1])
        print("Weights is sorted like: " + str(weights))
        vect[-1] = weights[-1][0]
        total_runs += 1
        print("Actual: " + str(vect[-1]))
        print("Correct: " + str(vect[2]))

        """And from here we will calculate error and such"""
        #This portion is used to calculate the error cost of the given case. Is added to the pool of average cost
        pure_weights = [None, None, None, None, None]
        correction_y = [None, None, None, None, None]
        for w in weights:
            pure_weights[w[0]] = w[1]
            #Indicates that the given weighting is the correct classification
            if w[0] == vect[2]:
                correction_y[w[0]] = 1
            else:
                correction_y[w[0]] = 0

        j = 0
        while j < len(WEIGHT_MATRIX):
            weight_change = []
            print("This is of value " + str(j) + ", but the correct val. is " + str(vect[2]))
            curr_y = 1 if j == vect[2] else 0

            k = 0
            while k < len(WEIGHT_MATRIX[j]):
                #change = vect[k] * dsigmoid(WEIGHT_MATRIX[j][k] * vect[k]) * 2 * (weights[j][1] - curr_y)
                change = vect[k] * (WEIGHT_MATRIX[j][k] * vect[k]) * 2 * (weights[j][1] - curr_y)
                print("Change element is " + str(change))
                print("Because of " + str(vect[k]) + " * (" + str(WEIGHT_MATRIX[j][k]) + " * " +  str(vect[k]) + ") * 2 * (" + str(weights[j][1]) + " - " + str(curr_y) + ")")
                weight_change.append(change) if not np.isnan(change) else weight_change.append(0.0)
                k += 1
            print("Weight change " + str(j) + " is: " + str(weight_change))
            change_matrix[j] += weight_change
            print("Change matrix is now " + str(change_matrix.shape))
            j += 1

        if vect[-1] == vect[2]:
            correct_runs += 1


        if vect[-1] == 0:
            head_size = vect[7] if head_size < vect[7] else head_size
    print("The algorithm is about " + str(100 * correct_runs/total_runs) + "% correct.")
    change_matrix = (change_matrix * STEP) / total_runs
    print(change_matrix)



    head = (0, [72, 1, 1, 0])
    subhead = (1, [45, 0, 1, 0])
    subsec = (2, [28, 0, 1, 0])
    body = (3, [28, 0, 0, 0])
    subtitle = (4, [20, 0, 0, 1])

    final_hier = [head, subhead, subsec, body, subtitle]

    return [correct_runs/total_runs, change_matrix]


def vector_fill_text(first_slide, run_external):
    """Fills out the np array of all text runs in the document
    Attributes of each vector are in the order:
    1)Run ID 2)Shape ID 3)"""
    """Keeps tabs on the id number of each run and retains them in an array"""
    id_shp = 0
    id_run = 0
    vector_arr = []
    del run_external[:]
    for shape in first_slide.shapes:
        if shape.has_text_frame:
            id_shp += 1
            for par in shape.text_frame.paragraphs:
                for run in par.runs:
                    #if run.font.name == 'Arial Black':
                        #run.font.bold = True
                    f_black = 0
                    f_size = 0
                    if (run.font.name is None) or (len(run.font.name) < 5):
                        f_black = 0
                    else:
                        f_black = int(run.font.name[-5:] in ("Black", "Heavy")) * (10000000)
                    if run.font.size is None:
                        f_size = 0
                    else:
                        f_size = run.font.size
                    f_bold = int(run.font.bold == True) * (5000000)
                    f_italic = int(2 // (int(run.font.italic == True) + 1) * (-100000))
                    x = shape.left // 10000
                    y = shape.top // 10000
                    width = shape.width // 100
                    height = shape.height // 100
                    l_space = par.line_spacing
                    level = par.level
                    length = int(200000 // (len(run.text) + 1))
                    run_attr = [id_run, id_shp, length, width, height, x, y, f_size, f_black, f_bold, f_italic, 0]
                    vector_arr.append(run_attr)
                    run_external.append(run)
                    id_run += 1
    nump_arr = np.array(vector_arr)
    return nump_arr


def find_centroids(nparray):
    vect_dim = nparray.shape[1]
    i = 0
    centroids = []

    while i < K:
        centroids.append(np.random.rand(1, vect_dim))
        # centroids.append(np.array([0, 0, rand.randint(1,150), rand.randint(20000, 500000), rand.randint(10000, 300000),
        #                            rand.randint(50, 3000), rand.randint(30, 3000), rand.randint(200000, 15000000),
        #                            rand.randint(0, 10000000), 0, 0, 0]))
        i += 1

    return centroids


def centroid_compare(nparray, centroids):
    for vect in nparray:
        min_clust_dist = None
        clust_ind = 0
        while clust_ind < len(centroids):
            if min_clust_dist is None:
                min_clust_dist = np.linalg.norm(centroids[clust_ind] - vect)
                vect[-1] = clust_ind
            elif min_clust_dist > np.linalg.norm(centroids[clust_ind] - vect):
                min_clust_dist = np.linalg.norm(centroids[clust_ind] - vect)
                vect[-1] = clust_ind
            clust_ind += 1
            #print(vect[-1])


def redef_centroids(nparray, centroids):
    clust_ind = 0
    new_centrs = []
    while clust_ind < len(centroids):
        new_centr = None
        clust_size = 0
        for vect in nparray:
            if vect[-1] == clust_ind:
                #print("passed in!")
                clust_size += 1
                if new_centr is None:
                    #print("OK")
                    new_centr = vect
                else:
                    new_centr = vect + new_centr
        if clust_size > 0:
            new_centrs.append(new_centr//clust_size)
            centroids[clust_ind] = list(new_centr//clust_size)
            print(centroids[clust_ind])
        else:
            print(centroids[clust_ind])
            new_centrs.append(centroids[clust_ind].flatten())
        clust_ind += 1
    return np.array(new_centrs)
    #print(centroids)
    #return old_centroids


def centroid_convergence(old, new):
    old_mat = np.array(old)
    new_mat = np.array(new)
    print("Old shape is " + str(old_mat.shape) + " and new is " + str(new_mat.shape))
    diff = np.subtract(old_mat, new_mat)
    print(np.linalg.norm(diff))
    return np.linalg.norm(diff)


def kmeans_text(nparray):
    iterations = 0

    centroids = find_centroids(nparray)
    old_centroids = np.array(centroids)

    conv_value = UPPERBOUND
    while conv_value > 1:
        centroid_compare(nparray, centroids)
        centroids = redef_centroids(nparray, centroids)
        conv_value = centroid_convergence(old_centroids, centroids)
        print("Convergence value is " + str(conv_value))
        old_centroids = np.array(centroids)
        iterations += 1

    print(nparray[0:])
    print(iterations)
    return centroids


def rud_hierarchy(centroids):
    """This is the only function, so far, that references a specific index in our vector space that's not an
    edge case"""
    heir = []
    i = 0
    while i < len(centroids):
        heir.append((i, np.linalg.norm(centroids[i])))
        i += 1

    heir.sort(key=lambda x: x[1])
    print(heir)
    header = (heir[3][0], [int(centroids[heir[3][0]][7]), 1, 1, 0])
    subhead = (heir[2][0], [(int(header[1][0]//1.618)) // 1.618, 0, 1, 0])
    body = (heir[1][0], [int(subhead[1][0]), 0, 0, 0])
    subtitl = (heir[0][0], [int(body[1][0]//1.618), 0, 0, 1])
    ret_array = [header, subhead, body, subtitl]
    return ret_array


def text_hier_alter(nparray, hier, runs):
    #Should have a companion function for font paring, likely implemented when giving the hier traits
    i = 0
    while i < K:
        for vect in nparray:
            if vect[-1] == hier[i][0]:
                print("I'm in")
                runs[vect[0]].font.size = Pt(hier[i][1][0])
                #Body check for bolding to ensure subheads and highlighting remains
                if not(i == 3 and runs[vect[0]].font.bold is True):
                    runs[vect[0]].font.bold = bool(hier[i][1][2])
                runs[vect[0]].font.italic = bool(hier[i][1][3])
                print(runs[vect[0]].text)
                print(vect[1])
        i += 1


def make_group_array(numpy_shape):
    group_array = []
    i = 0
    while i < len(numpy_shape):
        group_size = 0
        """Group Elements are in the order: Image, Chart, Table, Header, Subheader, Subsec, Body, Subtitle"""
        group_element = np.array([i, 0, 0, 0, 0, 0, 0, 0, 0])
        for vect in numpy_shapes:
            if vect[-1] == i:
                group_size += 1
                group_element[1:] += vect[6:-1]
        i += 1
        if group_size == 0:
            break
        group_array.append(group_element)
    return np.array(group_array)

# def group_parse(numpy_group, numpy_shape, nparray):
#     for group in numpy_group:
#         #Checks for body text first
#         if group[7] > 0:
#             #Check for headers first
#             if group[4] > 0:
#                 highest_element = None
#                 for shape in numpy_shape:
#                     if shape[-1] == group[0]:
#                         if highest_element is None:
#                             highest_element = shape
#                         else:
#                             if highest_element[4] < shape[4]:
#                                 highest_element = shape
#                 if highest_element[9] > 0:


#def cost(vect):



def quick_fill_shape(group_ind, numpy_shape, shapes):
    i = 0
    while i < len(numpy_shape):
        if numpy_shape[i][-1] == group_ind:
            shapes[i].fill.solid()
            shapes[i].fill.fore_color.rgb = RGBColor(57, 255, 20)
        i += 1


my_runs = []
#my_shapes = first.shapes
#my_shapes.append(second.shapes)

numpy_arr = vector_fill_clean(first, my_runs)
numpy2 = vector_fill_clean(second, my_runs)
numpy_arr = np.concatenate((numpy_arr, numpy2))

print("This is the second numpy array " + str(numpy2))

print(numpy_arr)
print(numpy_arr.shape)
print(numpy_arr.dtype)
print(my_runs[0].text)

head = (0, [72, 1, 1, 0])
subhead = (1, [45, 0, 1, 0])
subsec = (2, [28, 0, 1, 0])
body = (3, [28, 0, 0, 0])
subtitle = (4, [20, 0, 0, 1])

text_hier = [head, subhead, subsec, body, subtitle]


#centrs = kmeans_text(numpy_arr)
#text_hier = rud_hierarchy(centrs)
accuracy = 0
max_acc = 0
external_change = CHANGE_MATRIX

decoy = 1
print("This is the OG unedited matrix: " + str(WEIGHT_MATRIX))
iterations  = 0
accuracies = []
while iterations < 1:
    while accuracy < 0.7:
        og_acc = accuracy
        print("Weight before addition: " + str(WEIGHT_MATRIX))
        WEIGHT_MATRIX += external_change
        print("Weight after addition: " + str(WEIGHT_MATRIX))
        acc_and_weight = distr_assigning(numpy_arr)
        accuracy = acc_and_weight[0]
        external_change = acc_and_weight[1]
        print("This is the change being applied: " + str(external_change))
        print("This is the unedited weight matrix: " + str(WEIGHT_MATRIX))
        print("This is iteration number " + str(decoy))
        accuracies.append(accuracy)
        if og_acc >= accuracy:
            print("Some wacky taffy shit has occurred")
            break
        decoy += 1

    if max_acc < accuracy:
        max_acc = accuracy
    WEIGHT_MATRIX = np.random.rand(5, 18) * np.random.randint(1, 10000)
    iterations += 1

print("The final accuracy in ten iterations was: " + str(accuracies))
print("The final weightings are " + str(WEIGHT_MATRIX))
print("The numpy array was " + str(numpy_arr[10:20]))


# w_head = (0, sigmoid(np.dot([0, 0, 0, 1 /(vect[3] ** 2), 1, 1 / (vect[5] ** 2), 1 / (vect[6] ** 2), 1/(vect[7] ** 2), 1,
#                              1, 1, 1, 1, 1, 1, 1] * WEIGHT_MATRIX[0], vect)))
#         w_subhead = (1, sigmoid(np.dot([0, 0, 0, 1/(vect[3] ** 2), 1, 1/(vect[5] ** 2), 1/(vect[6] ** 2), 1/(vect[7] ** 2), 1,
#                                 1, 1, 1, 1, 1, 1, 1] * WEIGHT_MATRIX[1], vect)))
#         w_subsec = (2, sigmoid(np.dot([0, 0, 0, 1/(vect[3] ** 2), 1, 1/(vect[5] ** 2), 1, 1/(vect[7] ** 2), 1, 1, 1, 1, 1, 1, 1,
#                                1] * SUBSECTION, vect)))
#         w_body = (3, sigmoid(np.dot([0, 0, 0, 1, 1, 1, 1, 1, -90000/(vect[8] ** 2) + 0.0000001, 1, 1, 1, 1, 1, 1, 1] * BODY,
#                             vect)))
#         w_subtitle = (4, sigmoid(np.dot([0, 0, 0, 1, 1, 1/(vect[5] ** 2), 1, 1, 1/(vect[8] ** 2), 1, 1, 1, 1, 1, 1, 1] *
#                                 CAPTION, vect)))
#         weights = [w_head, w_subhead, w_subsec, w_body, w_subtitle]


#HEAD = text_hier[0]
#print(centrs)
#print(HEAD)

#numpy_shapes = vector_fill_shape(first, numpy_arr)
#print(numpy_shapes)

##text_hier_alter(numpy_arr, text_hier, my_runs)


#means = meanshift(numpy_shapes)
#print(means)

#arrange_meanshift(means, numpy_shapes)
#print(numpy_shapes)

#print(make_group_array(numpy_shapes))
#quick_fill_shape(0, numpy_shapes, my_shapes)
#print(cluster.mean_shift(shape_vector_distill(numpy_shapes)))


# centroidsy = find_centroids(numpy_arr)
# old_centroidsy = np.array(centroidsy)
# np.reshape(old_centroidsy, (4, 1, 12))
# print(old_centroidsy.shape)
# print(centroidsy)
# centroid_compare(numpy_arr, centroidsy)
#
#
# #print(numpy_arr[30:50])
#
#
# centroidsy = redef_centroids(numpy_arr, centroidsy)
# #print(old_centroidsy)
# #print(np.array(centroidsy))
# centroid_convergence(old_centroidsy, centroidsy)
#
# #print(centroidsy)
# #print(old_centroidsy)


prs1.save('test2.pptx')

