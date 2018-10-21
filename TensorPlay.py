from pptx import Presentation
from pptx.util import Inches, Pt
import numpy as np
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT as PP_ALIGN
import math
import sklearn
from sklearn import cluster
from pptx.dml.color import RGBColor, _NoneColor
from pptx.enum.dml import MSO_COLOR_TYPE
import tensorflow as tf
from tensorflow import keras
import Playgrounds as pg
from colormath.color_objects import LabColor, sRGBColor, HSLColor
from colormath.color_conversions import convert_color as cc
from copy import copy



def last_lett(word):
    i = -1
    while i > -len(word):
        if not word[i].isspace():
            return word[i]
        i -= 1

def rgb2hex(r,g,b):
    hex = "#{:02x}{:02x}{:02x}".format(r,g,b)
    return hex


def vector_fill_clean(first_slide, run_external):
    """Fills out the np array of all text runs in the document
    Attributes of each vector are in the order:
    1)Run ID 2)Shape ID 3)"""
    """Keeps tabs on the id number of each run and retains them in an array"""
    id_par = 0
    id_run = 0
    vector_arr = []
    id_arr = []
    avg_text_col = np.array([0, 0, 0])
    text_col_count = 0
    print(first_slide.follow_master_background)
    first_slide.background.fill.solid()
    first_slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)
    #first_slide.follow_master_background = False
    del run_external[:]
    for shape in first_slide.shapes:
        if shape.has_text_frame:
            shape.word_wrap = True
            shape.fill.background()
            shape.line.fill.background()
            id_shp = first_slide.shapes.index(shape)
            for par in shape.text_frame.paragraphs:
                #print(par.text)
                #print("This is paragraph: " + str(id_par))
                #print("*********")
                for run in par.runs:
                    print("*:STARTING:*", run.text)
                    correct_id = -1
                    if run.font.color.type is not None:
                        if run.font.color.type is MSO_COLOR_TYPE.RGB:
                            if run.font.color.rgb == RGBColor.from_string("FA6300"):
                                correct_id = 0
                            elif run.font.color.rgb == RGBColor.from_string("1398FF"):
                                correct_id = 1
                            elif run.font.color.rgb == RGBColor.from_string("13BA33"):
                                correct_id = 2
                            elif run.font.color.rgb == RGBColor.from_string("000000"):
                                correct_id = 3
                            elif run.font.color.rgb == RGBColor.from_string("EB1500"):
                                correct_id = 4
                            avg_text_col += (np.array(pg.rough_rgb_hex(str(run.font.color.rgb))) * len(run.text))
                            text_col_count += len(run.text)
                    if (run.font.name is None) or (len(run.font.name) < 5):
                        f_black = 0
                    else:
                        f_black = int(run.font.name[-5:] in ("Black", "Heavy"))

                    if run.font.size is None:
                        f_size = 0
                    else:
                        f_size = float(run.font.size.pt)
                    f_bold = int(run.font.bold == True)
                    f_italic = int(run.font.italic == True)
                    width = shape.width // 1000000
                    height = shape.height // 1000000
                    x = (shape.left // 1000000) + (width / 2)
                    y = (shape.top // 1000000) + (height / 2)
                    l_space = par.line_spacing if par.line_spacing is not None else 0
                    level = par.level
                    length = len(run.text)
                    punc = (int(last_lett(run.text) in (".", ",", "!", "?"))) if len(run.text) > 0 else 0
                    sect = int(run.text in ('Abstract', 'ABSTRACT', 'Introduction', 'INTRODUCTION', 'Method', 'METHOD',
                                            'Results', 'RESULTS', "Acknowledgments", "ACKNOWLEDGEMENTS", "Discussion",
                                            "DISCUSSION", "References", "REFERENCES", "Conclusions", "CONCLUSIONS"))
                    align = int(par.alignment == PP_ALIGN.CENTER)
                    run_attr = [id_run, id_shp, length, width, height, x, y, f_size, f_black, f_bold, f_italic, punc,
                                align, sect, l_space, level, 0]
                    # Only apply the if for the training phase
                    #if correct_id != -1:
                    vector_arr.append(run_attr)
                    correction_arr = [0, 0, 0, 0, 0]
                    correction_arr[correct_id] = 1
                    id_arr.append(correction_arr)
                    run_external.append(run)
                    id_run += 1
                id_par += 1
    nump_arr = np.array(vector_arr)
    correct_arr = np.array(id_arr)
    new_avg = np.array(avg_text_col) // text_col_count
    #print(nump_arr[40:])
    print("Hey the average RGB is " + str(new_avg))
    return [nump_arr, correct_arr, new_avg]


def vector_fill_clean_training(first_slide, run_external):
    """Fills out the np array of all text runs in the document
    Attributes of each vector are in the order:
    1)Run ID 2)Shape ID 3)"""
    """Keeps tabs on the id number of each run and retains them in an array"""
    id_par = 0
    id_run = 0
    vector_arr = []
    id_arr = []
    avg_text_col = np.array([0, 0, 0])
    text_col_count = 0
    first_slide.background.fill.solid()
    first_slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)
    del run_external[:]
    for shape in first_slide.shapes:
        if shape.has_text_frame:
            shape.word_wrap = True
            shape.fill.background()
            shape.line.fill.background()
            id_shp = first_slide.shapes.index(shape)
            for par in shape.text_frame.paragraphs:
                #print(par.text)
                #print("This is paragraph: " + str(id_par))
                #print("*********")
                for run in par.runs:
                    print("*:STARTING:*", run.text)
                    correct_id = -1
                    if run.font.color.type is not None:
                        if run.font.color.type is MSO_COLOR_TYPE.RGB:
                            if run.font.color.rgb == RGBColor.from_string("FA6300"):
                                correct_id = 0
                            elif run.font.color.rgb == RGBColor.from_string("1398FF"):
                                correct_id = 1
                            elif run.font.color.rgb == RGBColor.from_string("13BA33"):
                                correct_id = 2
                            elif run.font.color.rgb == RGBColor.from_string("000000"):
                                correct_id = 3
                            elif run.font.color.rgb == RGBColor.from_string("EB1500"):
                                correct_id = 4
                            avg_text_col += (np.array(pg.rough_rgb_hex(str(run.font.color.rgb))) * len(run.text))
                            text_col_count += len(run.text)
                    #print("The real answer is " + str(correct_id))
                    if (run.font.name is None) or (len(run.font.name) < 5):
                        f_black = 0
                    else:
                        f_black = int(run.font.name[-5:] in ("Black", "Heavy"))

                    if run.font.size is None:
                        f_size = 0
                    else:
                        f_size = float(run.font.size.pt)
                    f_bold = int(run.font.bold == True)
                    f_italic = int(run.font.italic == True)
                    width = shape.width // 1000000
                    height = shape.height // 1000000
                    x = (shape.left // 1000000) + (width / 2)
                    y = (shape.top // 1000000) + (height / 2)
                    l_space = par.line_spacing if par.line_spacing is not None else 0
                    level = par.level
                    length = len(run.text)
                    punc = (int(last_lett(run.text) in (".", ",", "!", "?"))) if len(run.text) > 0 else 0
                    sect = int(run.text in ('Abstract', 'ABSTRACT', 'Introduction', 'INTRODUCTION', 'Method', 'METHOD',
                                            'Results', 'RESULTS', "Acknowledgments", "ACKNOWLEDGEMENTS", "Discussion",
                                            "DISCUSSION", "References", "REFERENCES", "Conclusions", "CONCLUSIONS"))
                    align = int(par.alignment == PP_ALIGN.CENTER)
                    run_attr = [id_run, id_shp, length, width, height, x, y, f_size, f_black, f_bold, f_italic, punc,
                                align, sect, l_space, level, 0]
                    # Only apply the if for the training phase
                    if correct_id != -1:
                        vector_arr.append(run_attr)
                        correction_arr = [0, 0, 0, 0, 0]
                        correction_arr[correct_id] = 1
                        id_arr.append(correction_arr)
                        run_external.append(run)
                        id_run += 1
                id_par += 1
    nump_arr = np.array(vector_arr)
    correct_arr = np.array(id_arr)
    new_avg = np.array(avg_text_col) // text_col_count
    #print(nump_arr[40:])
    print("Hey the average RGB is " + str(new_avg))
    return [nump_arr, correct_arr]


def pptx_data(filename, external_runs):
    """Returns in the format <features>, <labels>"""
    prs = Presentation(filename)
    slide = prs.slides[0]

    return vector_fill_clean(slide, external_runs)[0], prs, vector_fill_clean(slide, external_runs)[2]


"""And down bellow begins the most tense flows that've ever graced your property"""


def color_scheme_translate(avg_col_arr):
    """Takes in the average color of a given function, and spits back out the monochromatic color scheme along said
    path. Header is made a complimentary color to the subheader."""
    avg_col = sRGBColor(avg_col_arr[0], avg_col_arr[1], avg_col_arr[2])
    avg_hsl_source = cc(avg_col, HSLColor)
    hsl_base = copy(avg_hsl_source)
    hsl_subsec = copy(avg_hsl_source)
    hsl_subhead = copy(avg_hsl_source)
    hsl_base.hsl_l = 20
    hsl_subsec.hsl_l = 40
    hsl_subhead.hsl_l = 60
    hsl_head = pg.compliment_rgb(copy(hsl_subhead))

    print(str(hsl_head), "vs", str(hsl_base))


    # Converts the colors back to their RGB form, and reformats them to be compatible with pptx

    avbase = cc(hsl_base, sRGBColor)
    avbase_hex = rgb2hex(int((avbase.rgb_r)), int((avbase.rgb_g)), int((avbase.rgb_b)))
    avg_base = RGBColor.from_string(avbase_hex[1:])

    avbase1 = cc(hsl_subhead, sRGBColor)
    avbase1_hex = rgb2hex(int((avbase1.rgb_r)), int((avbase1.rgb_g)), int((avbase1.rgb_b)))
    avg_subhead = RGBColor.from_string(avbase1_hex[1:])

    avbase2 = cc(hsl_subsec, sRGBColor)
    avbase2_hex = rgb2hex(int((avbase2.rgb_r)), int((avbase2.rgb_g)), int((avbase2.rgb_b)))
    avg_subsec = RGBColor.from_string(avbase2_hex[1:])

    avbase3 = cc(hsl_head, sRGBColor)
    avbase3_hex = rgb2hex(int((avbase3.rgb_r)), int((avbase3.rgb_g)), int((avbase3.rgb_b)))
    avg_head = RGBColor.from_string(avbase3_hex[1:])

    return [avg_head, avg_subhead, avg_subsec, avg_base]


def main():
    prs_set = ['thissatest copy.pptx',  'leepostertest copy.pptx', 'symposiumPoster2.pptx',
               'materialstailgate11132015.pptx']

    external_runs1 = []
    #features = np.array([[], []])
    #labels = np.array([[], []])

    prs1 = Presentation('icassp2017_poster.pptx')
    slide1 = prs1.slides[0]
    features, labels = vector_fill_clean_training(slide1, external_runs1)

    for name in prs_set:
        prs = Presentation(name)
        slide = prs.slides[0]
        feats, labs = vector_fill_clean_training(slide, external_runs1)
        print(labs.shape)
        features = np.concatenate((features, feats))
        labels = np.concatenate((labels, labs))



    #prs2 = Presentation('thissatest copy.pptx')
    #prs1 = Presentation('sectest.pptx')
    #prs1 = Presentation('leepostertest copy.pptx')
    #first = prs1.slides[0]
    #second = prs2.slides[0]
    #title = first.shapes[0]
    #external_runs = []
    #external_runs2 = []

    node_count_l1 = 17
    class_count = 5
    batch_size = 16

    #features, labels, arb1 = vector_fill_clean(first, external_runs)
    #features2, labels2, arb2 = vector_fill_clean(second, external_runs2)
    #features = np.concatenate((features, features2))
    #labels = np.concatenate((labels, labels2))
    #external_runs.append(external_runs2)

    assert features.shape[0] == labels.shape[0]

    model = keras.Sequential()

    model.add(keras.layers.Dense(node_count_l1, input_dim=17, activation='relu'))
    model.add(keras.layers.Dense(class_count, activation='softmax'))

    model.compile(optimizer=tf.keras.optimizers.Adam(lr=0.001), loss='categorical_crossentropy', metrics=['accuracy'])

    model.fit(features, labels, epochs=100, batch_size=batch_size)

    # Saving the current model below
    model.save('trained_text_model.h5')


#main()



